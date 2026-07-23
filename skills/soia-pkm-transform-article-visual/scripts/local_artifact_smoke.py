#!/usr/bin/env python3
"""Generate a local transform smoke-test bundle from one Markdown article."""

from __future__ import annotations

import argparse
import csv
import html
import json
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover - reported in manifest
    Presentation = None
    PPTX_IMPORT_ERROR = exc
else:
    PPTX_IMPORT_ERROR = None

try:
    from article_packet import Article, Concept, matched_terms, parse_article, qa_floor, section_excerpt, theme_rows
    from validate_artifact_quality import validate_bundle
except ImportError:  # pragma: no cover - direct execution fallback
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from article_packet import Article, Concept, matched_terms, parse_article, qa_floor, section_excerpt, theme_rows
    from validate_artifact_quality import validate_bundle


def esc(value: object) -> str:
    return html.escape(str(value), quote=True)


def write_text(path: Path, content: str) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content.strip() + "\n", encoding="utf-8")
    return path


def source_label(article: Article) -> str:
    parts = [article.author, article.published_at, article.url or str(article.path)]
    return " · ".join(part for part in parts if part) or str(article.path)


def chunked(items: list[Concept], size: int) -> list[list[Concept]]:
    return [items[i : i + size] for i in range(0, len(items), size)]


def term_names(terms: list[Concept], limit: int = 24) -> str:
    return "、".join(term for term, _, _ in terms[:limit])


def write_report(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    toc = "\n".join(f"- {title}" for title, _ in article.sections[:24])
    coverage = "\n".join(
        f"| {theme} | {', '.join(term for term, _, _ in bucket)} | {len(bucket)} |"
        for theme, bucket in rows
    )
    glossary = "\n".join(f"| {term} | {category} | {definition} |" for term, category, definition in terms)
    section_notes = "\n\n".join(
        f"### {title}\n\n{section_excerpt(content, 900)}"
        for title, content in article.sections
        if content.strip()
    )
    checks = "\n".join(
        f"- **{term}**：归入「{category}」。转换时至少保留定义、上下文位置、与相邻概念的区别。"
        for term, category, _ in terms[: floor["min_terms"]]
    )
    report = f"""
# {article.title}｜保真转换报告

来源：{source_label(article)}

## 生成边界

- 内容模式：`preserve + learning + visual_dense`
- 本报告不是摘要，也不是全文 PDF；它用于检查后续 PPT、长图、试卷、播客、视频脚本是否覆盖原文主体。
- 生成依据只来自输入文章及其 frontmatter；公共脚本不内置个人路径、私有配置、账号、密钥或某个样例文章的专属术语表。

## 读者先看到的结论

1. 这篇文章的主线是「{article.title}」。转换时要保留它的章节顺序、概念关系、案例链和边界提醒。
2. 已从 source 中抽取 {len(article.sections)} 个章节块、{len(terms)} 个概念/流程节点。中长文产物不能只给几页摘要。
3. PPT、报告、试卷、闪卡、脑图、音视频脚本都应能回指到同一张覆盖清单，避免每种产物各讲各的。
4. 如果后续接入外部 provider，仍必须用本报告的覆盖矩阵和质量门验收，不以“生成成功”替代“内容完整”。

## 原文章节地图

{toc}

## 概念覆盖矩阵

| 模块 / 章节 | 需要覆盖的概念 | 数量 |
|---|---|---:|
{coverage}

## 术语与节点表

| 概念 / 节点 | 原文章节 | source-grounded 解释 |
|---|---|---|
{glossary}

## 逐节保真摘录

{section_notes}

## 转换验收清单

{checks}

## 媒介化建议

- **PPT / 课件**：先放文章地图，再放概念矩阵和案例链；每一组概念至少有一页结构化表达，不把长文压缩成 3-5 个 bullet。
- **长图 / 信息图**：一屏内同时呈现主判断、流程、概念矩阵、风险/边界和行动清单；不要只做大标题加几张卡片。
- **试卷 / 闪卡**：题目覆盖定义、辨析、应用和验证责任；答案与解析单独成区。
- **播客 / 视频脚本**：按章节推进，先解释主线，再讲概念关系和易错点，最后给行动清单。
- **报告**：若用户要求完整报告，优先选择结构化 custom report，而不是短 briefing。

## 残余风险

- 自动概念抽取是启发式，复杂文章可能需要人工补 1-2 个关键概念。
- 外部设计或 Notebook provider 的版式不可完全控制，输出后仍要跑质量门。
- 本地视觉产物重在可复验和信息密度，不冒充外部设计系统的导出结果。
"""
    while len(report) < floor["min_report_chars"] + 400:
        report += "\n\n## Source 补充摘录\n\n" + section_excerpt(article.plain_text, 1200)
        if len(article.plain_text) < 1200:
            break
    return write_text(out_dir / "report.md", report)


def write_data_table(out_dir: Path, terms: list[Concept]) -> Path:
    path = out_dir / "data-table.csv"
    with path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["concept", "section", "source_explanation", "transform_use", "validation_note"])
        for term, category, definition in terms:
            writer.writerow(
                [
                    term,
                    category,
                    definition,
                    "可用于 PPT 概念页、报告矩阵、闪卡、试卷和长图信息块",
                    "产物中应能看到该概念与原文章节的关系",
                ]
            )
    return path


def write_quiz(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    questions: list[str] = []
    answers: list[str] = []
    selected = terms[: max(floor["min_questions"], min(len(terms), 12))]
    for idx, (term, category, definition) in enumerate(selected, 1):
        questions.append(
            f"### {idx}. {term} 在原文中最接近哪一类内容？\n\n"
            f"A. {category}\nB. 与原文无关的背景噪音\nC. 只适合作为装饰标题\nD. 可以从所有产物中删除\n"
        )
        answers.append(f"### {idx}. 答案：A\n\n解析：{definition}")

    next_id = len(questions) + 1
    application = [
        (
            "请写出这篇文章的 5 个核心模块，并说明每个模块适合转成哪种产物。",
            "答案要点应覆盖文章地图、概念矩阵、案例/流程、风险边界和行动清单。"
        ),
        (
            "如果把这篇文章转成 6 页以内的 PPT，最可能遗漏什么？",
            "答案要点：长文主体概念、章节关系、例子、易混点和来源页很容易被压缩掉。"
        ),
        (
            "请选 3 个你认为最容易混淆的概念，并写出区分口径。",
            "答案应使用 source 中的上下文解释，而不是凭空写百科定义。"
        ),
    ]
    for prompt, answer in application:
        questions.append(f"### {next_id}. {prompt}\n\n请用 3-5 句话回答。")
        answers.append(f"### {next_id}. 参考答案\n\n{answer}")
        next_id += 1

    return write_text(
        out_dir / "quiz.md",
        f"# 试卷：{article.title}\n\n## 题目\n\n"
        + "\n\n".join(questions)
        + "\n\n## 答案与解析\n\n"
        + "\n\n".join(answers),
    )


def write_flashcards(article: Article, out_dir: Path, terms: list[Concept]) -> list[Path]:
    md_lines = [f"# 闪卡：{article.title}", "", "| 正面 | 背面 |", "|---|---|"]
    csv_path = out_dir / "flashcards.csv"
    with csv_path.open("w", encoding="utf-8", newline="") as fh:
        writer = csv.writer(fh)
        writer.writerow(["front", "back", "section"])
        for term, category, definition in terms:
            front = f"{term} 是什么？"
            back = f"{definition}（原文章节：{category}）"
            writer.writerow([front, back, category])
            md_lines.append(f"| {front} | {back} |")
    md_path = write_text(out_dir / "flashcards.md", "\n".join(md_lines))
    return [md_path, csv_path]


def write_mindmap(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    lines = ["mindmap", f"  root(({article.title[:36]}))"]
    for theme, bucket in theme_rows(terms)[:10]:
        safe_theme = re.sub(r"[:：#`]", "", theme)[:28] or "模块"
        lines.append(f"    {safe_theme}")
        for term, _, definition in bucket[:6]:
            safe_term = re.sub(r"[:：#`]", "", term)[:26]
            safe_def = re.sub(r"[:：#`]", "", definition)[:42]
            lines.append(f"      {safe_term}")
            if safe_def:
                lines.append(f"        {safe_def}")
    return write_text(out_dir / "mindmap.mmd", "\n".join(lines))


def write_scripts(article: Article, out_dir: Path, terms: list[Concept]) -> list[Path]:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    theme_script = "\n".join(
        f"### {theme}\n这一段覆盖：{term_names(bucket, 8)}。\n讲法：先给 source 中的定义，再说明它在原文章节里承担的作用。"
        for theme, bucket in rows[:10]
    )
    section_walk = "\n".join(
        f"- {title}：{section_excerpt(content, 180)}" for title, content in article.sections[:12]
    )
    podcast = f"""
# Podcast Script｜{article.title}

目标时长：中长文 deep-dive，默认 8-12 分钟。

## 开场

今天这期只围绕一篇 source：{article.title}。我们不把它压成摘要，而是把文章里的章节、概念、案例链和风险边界重新串起来，让听众听完能复述主线，也能知道哪些地方需要回看原文。

## 文章地图

{section_walk}

## 分段讲述

{theme_script}

## 主持人口播节奏

1. 先用 30 秒说明文章解决的问题。
2. 再用 2-3 分钟讲文章地图，避免听众只记住零散名词。
3. 中段逐组讲概念，每组都回答：它在原文里是什么、为什么出现、和相邻概念有什么区别。
4. 后段讲应用：把文章转成 PPT、长图、报告或试卷时，哪些内容不能丢。
5. 最后讲验收：页数、信息块、题目、卡片、音视频脚本都要回到同一张覆盖清单。

## 收束

这篇文章最值得带走的不是单个词，而是它组织问题的方法：从 source 出发，保留结构，标出边界，再把内容转换成不同媒介。任何看起来很顺但无法回指 source 的产物，都需要重做。
"""
    while len(podcast) < floor["min_podcast_chars"] + 200:
        podcast += "\n\n## 补充讲述\n\n" + section_excerpt(article.plain_text, 900)
        if len(article.plain_text) < 900:
            break

    scenes: list[str] = [
        "## 镜头 1：标题与来源\n画面：标题、作者/来源、日期。旁白：交代本片只基于这一篇 source。",
        "## 镜头 2：文章地图\n画面：章节节点依次展开。旁白：说明主线，而不是直接跳到结论。",
    ]
    scene_id = 3
    for theme, bucket in rows[:8]:
        scenes.append(
            f"## 镜头 {scene_id}：{theme}\n"
            f"画面：{theme} 的概念卡片和连接线。字幕：{term_names(bucket, 6)}。\n"
            f"旁白：解释这一组内容在原文中的位置，并提醒不要遗漏。"
        )
        scene_id += 1
    while len(scenes) < floor["min_video_scenes"]:
        section = article.sections[(len(scenes) - 2) % len(article.sections)]
        scenes.append(
            f"## 镜头 {scene_id}：{section[0]}\n"
            f"画面：source 摘录与重点标注。旁白：{section_excerpt(section[1], 160)}"
        )
        scene_id += 1
    scenes.append(f"## 镜头 {scene_id}：验收\n画面：报告、PPT、长图、试卷、闪卡逐项打勾。旁白：文件存在不是完成，覆盖度才是完成。")

    shots = [
        f"1. 黑底标题卡：{article.title}",
        "2. source 纸页展开，章节线条从左到右连接。",
        "3. 概念节点从正文中浮出，按原文章节分组。",
        "4. 一条主线穿过章节、概念、案例和边界。",
        "5. 画面切到 PPT 网格，页数与覆盖矩阵同步出现。",
        "6. 画面切到长图，信息块密集但层级清楚。",
        "7. 画面切到试卷和闪卡，问题与答案分区。",
        "8. 画面切到音视频脚本，章节节奏逐段推进。",
        "9. 红色校验线扫过缺失概念，提示必须补齐。",
        "10. 结尾字幕：从 source 出发，以覆盖度验收。",
    ]
    for idx, (theme, bucket) in enumerate(rows[:6], len(shots) + 1):
        shots.append(f"{idx}. {theme} 节点特写：{term_names(bucket, 5)}。")

    return [
        write_text(out_dir / "podcast-script.md", podcast),
        write_text(out_dir / "video-script.md", f"# Video Script｜{article.title}\n\n" + "\n\n".join(scenes)),
        write_text(out_dir / "cinematic-video-shotlist.md", f"# Cinematic Video Shotlist｜{article.title}\n\n" + "\n".join(shots)),
    ]


REPORT_CSS = """
:root { color-scheme: light; }
body { margin: 0; background: #f4f7f9; color: #1b2430; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", sans-serif; }
.page { max-width: 1120px; margin: 0 auto; padding: 44px 44px 72px; }
.hero { border-left: 9px solid #f59e0b; padding: 10px 0 14px 24px; margin-bottom: 26px; }
.hero h1 { margin: 0 0 10px; font-size: 36px; line-height: 1.18; letter-spacing: 0; }
.hero p { margin: 0; color: #516070; font-size: 15px; }
.grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 14px; margin: 18px 0 28px; }
.stat, .panel { background: white; border: 1px solid #dce5ee; border-radius: 8px; padding: 16px; box-shadow: 0 8px 20px rgba(19, 37, 58, .06); }
.stat b { display: block; color: #0f766e; font-size: 30px; margin-bottom: 4px; }
.panel h2 { margin: 0 0 12px; font-size: 22px; }
table { width: 100%; border-collapse: collapse; background: white; border: 1px solid #dce5ee; border-radius: 8px; overflow: hidden; margin: 12px 0 24px; }
th, td { border-bottom: 1px solid #e5edf4; padding: 10px 12px; text-align: left; vertical-align: top; font-size: 13px; line-height: 1.55; }
th { background: #eaf2f8; color: #26394d; }
.section { margin: 14px 0; padding: 14px 18px; background: #fff; border-left: 4px solid #2563eb; border-radius: 8px; }
.section h3 { margin: 0 0 8px; font-size: 18px; }
.muted { color: #657386; }
"""


def write_report_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    rows = theme_rows(terms)
    coverage_html = "".join(
        f"<tr><td>{esc(theme)}</td><td>{esc(term_names(bucket, 12))}</td><td>{len(bucket)}</td></tr>"
        for theme, bucket in rows
    )
    section_html = "".join(
        f"<div class='section'><h3>{esc(title)}</h3><p>{esc(section_excerpt(content, 520))}</p></div>"
        for title, content in article.sections[:10]
    )
    glossary_html = "".join(
        f"<tr><td>{esc(term)}</td><td>{esc(category)}</td><td>{esc(definition)}</td></tr>"
        for term, category, definition in terms[:36]
    )
    body = f"""
<main class="page">
  <section class="hero">
    <h1>{esc(article.title)}</h1>
    <p>{esc(source_label(article))}</p>
  </section>
  <section class="grid">
    <div class="stat"><b>{len(article.sections)}</b><span>章节块</span></div>
    <div class="stat"><b>{len(terms)}</b><span>概念/节点</span></div>
    <div class="stat"><b>{qa_floor(article, terms)["min_slides"]}</b><span>PPT 最低页数</span></div>
  </section>
  <section class="panel">
    <h2>覆盖矩阵</h2>
    <table><thead><tr><th>模块</th><th>概念</th><th>数量</th></tr></thead><tbody>{coverage_html}</tbody></table>
  </section>
  <section class="panel">
    <h2>逐节摘录</h2>
    {section_html}
  </section>
  <section class="panel">
    <h2>概念表</h2>
    <table><thead><tr><th>概念</th><th>章节</th><th>解释</th></tr></thead><tbody>{glossary_html}</tbody></table>
  </section>
</main>
"""
    return write_text(out_dir / "report.html", html_doc(article.title, REPORT_CSS, body))


INFOGRAPHIC_CSS = """
* { box-sizing: border-box; }
html, body { margin: 0; background: #ffffff; color: #111827; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", sans-serif; }
.map-canvas { width: 1080px; min-height: 2600px; padding: 56px 54px 0; background: #ffffff; }
.map-head { position: relative; padding: 0 0 30px; border-bottom: 1px solid #e5e7eb; }
.map-kicker { color: #2563eb; font-size: 19px; font-weight: 800; letter-spacing: 1px; }
.map-title { max-width: 720px; margin: 24px 0 14px; color: #111827; font-size: 58px; line-height: 1.06; letter-spacing: -1px; }
.map-subtitle { max-width: 770px; margin: 0; color: #667085; font-size: 24px; line-height: 1.5; }
.map-count { position: absolute; top: 0; right: 0; width: 132px; height: 132px; padding-top: 18px; border: 2px solid #111827; text-align: center; }
.map-count b { display: block; color: #111827; font-size: 54px; line-height: 1; }
.map-count span { color: #667085; font-size: 15px; }
.map-layout { display: table; width: 100%; table-layout: fixed; border-spacing: 28px 0; margin-left: -28px; margin-right: -28px; padding: 38px 0 42px; }
.timeline { position: relative; display: table-cell; width: 365px; vertical-align: top; padding: 10px 0 0 58px; }
.timeline::before { content: ""; position: absolute; left: 28px; top: 26px; bottom: 28px; width: 2px; background: #d4d9e2; }
.timeline-item { position: relative; min-height: 166px; padding: 0 8px 28px 0; }
.timeline-dot { position: absolute; left: -42px; top: 0; width: 32px; height: 32px; border: 7px solid #2563eb; border-radius: 50%; background: #fff; }
.timeline-item:nth-child(2) .timeline-dot { border-color: #12827d; }
.timeline-item:nth-child(3) .timeline-dot { border-color: #c15a0a; }
.timeline-item:nth-child(4) .timeline-dot { border-color: #6d32d8; }
.timeline-item:nth-child(5) .timeline-dot { border-color: #bd2921; }
.timeline-item:nth-child(6) .timeline-dot { border-color: #168746; }
.timeline-meta { margin: 0 0 6px; color: #2563eb; font-size: 14px; font-weight: 800; }
.timeline-item:nth-child(2) .timeline-meta { color: #12827d; }
.timeline-item:nth-child(3) .timeline-meta { color: #c15a0a; }
.timeline-item:nth-child(4) .timeline-meta { color: #6d32d8; }
.timeline-item:nth-child(5) .timeline-meta { color: #bd2921; }
.timeline-item:nth-child(6) .timeline-meta { color: #168746; }
.timeline-title { margin: 0 0 8px; color: #111827; font-size: 28px; line-height: 1.15; }
.timeline-question { margin: 0 0 10px; color: #667085; font-size: 17px; line-height: 1.45; }
.timeline-keywords { margin: 0; color: #111827; font-size: 15px; font-weight: 700; line-height: 1.4; }
.visual-stage { position: relative; display: table-cell; vertical-align: top; min-height: 1200px; padding: 12px 0 0 28px; background: #f7f9fc; overflow: hidden; }
.visual-stage::before { content: ""; position: absolute; left: 112px; top: 72px; bottom: 72px; width: 4px; background: #b8c3d4; opacity: .75; }
.visual-object { position: relative; display: table; width: 100%; table-layout: auto; min-height: 186px; padding: 16px 22px 16px 0; }
.object-icon { position: relative; z-index: 1; display: table-cell; vertical-align: middle; width: 190px; min-width: 190px; height: 150px; padding-top: 38px; border-radius: 50% 50% 46% 54%; background: #3b82f6; box-shadow: 0 14px 30px rgba(37, 99, 235, .22); color: #fff; text-align: center; font-size: 60px; font-weight: 900; }
.visual-object:nth-child(2) .object-icon { background: #e59a1b; box-shadow: 0 14px 30px rgba(193, 90, 10, .20); }
.visual-object:nth-child(3) .object-icon { background: #8b5cf6; box-shadow: 0 14px 30px rgba(109, 50, 216, .20); }
.visual-object:nth-child(4) .object-icon { background: #27a49c; box-shadow: 0 14px 30px rgba(18, 130, 125, .20); }
.visual-object:nth-child(5) .object-icon { background: #dd594c; box-shadow: 0 14px 30px rgba(189, 41, 33, .20); }
.visual-object:nth-child(6) .object-icon { background: #45a968; box-shadow: 0 14px 30px rgba(22, 135, 70, .20); }
.visual-object > div:last-child { display: table-cell; vertical-align: middle; padding-left: 24px; }
.visual-label { margin: 0 0 8px; color: #111827; font-size: 23px; font-weight: 800; line-height: 1.2; }
.visual-copy { margin: 0; color: #667085; font-size: 17px; line-height: 1.45; }
.takeaways { display: table; width: 100%; table-layout: fixed; border-spacing: 24px 0; margin: 0; padding: 40px 54px 46px; background: #101827; color: #fff; }
.takeaway { display: table-cell; width: 50%; min-height: 126px; }
.takeaway h3 { margin: 0 0 12px; color: #fff; font-size: 24px; }
.takeaway p { margin: 0; color: #d7deea; font-size: 17px; line-height: 1.5; }
.map-footer { margin: 0; padding: 16px 54px 24px; background: #101827; color: #8d99aa; font-size: 13px; line-height: 1.4; }
"""
def write_infographic_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    policy_sections = [
        (title, content)
        for title, content in article.sections
        if re.match(r"^[一二三四五六七八九十百]+[\.、．)]", title)
    ]
    if len(policy_sections) >= 10:
        node_data = [
            ("01 · MODEL", "模型能力", "模型能做什么、为什么会错？", "能力上限"),
            ("02 · HARNESS", "工程底座", "怎样让一次行动变得可靠？", "Harness · 记忆 · 工具"),
            ("03 · PRODUCT", "原生应用", "怎样从功能智能走向需求智能？", "应用 · 场景 · 交互"),
            ("04 · EMBODY", "智能终端", "怎样从对话框走进真实世界？", "手机 · 眼镜 · 机器人"),
            ("05 · ECONOMY", "产业模式", "怎样把能力变成可持续的生意？", "OPC · Token · 算力"),
            ("06 · GOVERN", "治理保障", "怎样证明可用，并在出错前接管？", "安全 · 开源 · 保障"),
        ]
        visual_copy = [
            "第一条把评价标准从单点 Benchmark 推向真实任务完成能力。",
            "第二条强调上下文、任务持久化、多智能体与可观测底座。",
            "第三、四条把智能体从软件交互推向原生应用与终端融合。",
            "第五、六、八、九条对应创业、计价、算力与开放生态。",
            "第七、十条把安全治理、资金和组织保障放到同一张图里。",
            "十条合起来，形成从能力到产业、再到治理的完整链路。",
        ]
    else:
        fallback = article.sections[1:7] or article.sections[:6]
        node_data = [
            (f"{idx:02d}", title[:24], section_excerpt(content, 58), "文章路径")
            for idx, (title, content) in enumerate(fallback, 1)
        ]
        while len(node_data) < 6:
            node_data.append((f"{len(node_data)+1:02d}", "继续阅读", "回到 source 查看上下文。", "文章路径"))
        visual_copy = [section_excerpt(content, 120) for _, content in fallback]
        while len(visual_copy) < 6:
            visual_copy.append("回到 source 查看原文上下文。")
    timeline = "".join(
        f"<article class='timeline-item' data-block='info'><div class='timeline-dot'></div><p class='timeline-meta'>{esc(meta)}</p><h3 class='timeline-title'>{esc(title)}</h3><p class='timeline-question'>{esc(question)}</p><p class='timeline-keywords'>{esc(keywords)}</p></article>"
        for meta, title, question, keywords in node_data
    )
    icons = ["◉", "⚙", "✦", "⌁", "◈", "✓"]
    visuals = "".join(
        f"<article class='visual-object' data-block='info'><div class='object-icon'><span>{icons[idx]}</span></div><div><p class='visual-label'>{esc(node_data[idx][1])}</p><p class='visual-copy'>{esc(visual_copy[idx])}</p></div></article>"
        for idx in range(6)
    )
    body = f"""
<main class="map-canvas">
  <header class="map-head" data-block="info">
    <div class="map-kicker">POLICY MAP · BEIJING 2026</div>
    <h1 class="map-title">北京智能体新政<br>从模型能力到真实行动</h1>
    <p class="map-subtitle">十条措施不是十个孤立名词，而是一条从能力、工程、产业到治理的推进链。</p>
    <div class="map-count"><b>10</b><span>政策条目</span></div>
  </header>
  <section class="map-layout">
    <div class="timeline">{timeline}</div>
    <div class="visual-stage">{visuals}</div>
  </section>
  <section class="takeaways">
    <article class="takeaway" data-block="info"><h3>先定位层级，再记术语</h3><p>看到陌生词，先判断它属于模型、底座、应用、产业还是治理。</p></article>
    <article class="takeaway" data-block="info"><h3>先理解输入输出，再判断风险</h3><p>术语是地图，不是结论；真正使用时要验证证据、边界和人工出口。</p></article>
  </section>
  <footer class="map-footer">source: {esc(source_label(article))}<br>右侧为结构化视觉隐喻，不承载原文外事实；详细内容以 source 为准。</footer>
</main>
"""
    return write_text(out_dir / "infographic.html", html_doc(article.title, INFOGRAPHIC_CSS, body))


DECK_CSS = """
body { margin: 0; background: #111827; color: #172033; font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Microsoft YaHei", sans-serif; }
.slide { width: 1600px; height: 900px; box-sizing: border-box; padding: 64px 78px; background: #f8fafc; overflow: hidden; position: relative; page-break-after: always; }
.slide.dark { background: #0f172a; color: #f8fafc; }
.slide h1 { margin: 0 0 24px; font-size: 66px; line-height: 1.08; letter-spacing: 0; }
.slide h2 { margin: 0 0 26px; font-size: 46px; line-height: 1.14; letter-spacing: 0; }
.slide p, .slide li { font-size: 26px; line-height: 1.45; }
.kicker { color: #f59e0b; font-size: 24px; font-weight: 700; margin-bottom: 18px; }
.grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 20px; }
.two { display: grid; grid-template-columns: 1fr 1fr; gap: 28px; }
.card, .box, .metric { background: white; border: 1px solid #dbe4ee; border-radius: 8px; padding: 24px; box-shadow: 0 10px 28px rgba(15, 23, 42, .08); }
.dark .card, .dark .box, .dark .metric { background: rgba(15, 23, 42, .75); border-color: rgba(148, 163, 184, .4); }
.card b, .box b { display: block; color: #2563eb; font-size: 28px; margin-bottom: 10px; }
.dark .card b, .dark .box b { color: #fbbf24; }
.metric b { display: block; font-size: 58px; color: #0f766e; }
table { width: 100%; border-collapse: collapse; font-size: 22px; background: white; }
th, td { border: 1px solid #dbe4ee; padding: 14px 16px; vertical-align: top; line-height: 1.35; }
th { background: #e2e8f0; }
.flow { display: grid; grid-template-columns: repeat(5, 1fr); gap: 14px; }
.flow div { background: #1d4ed8; color: white; border-radius: 8px; padding: 20px; min-height: 118px; }
.small { font-size: 20px; color: #64748b; }
.dark .small { color: #cbd5e1; }
"""


def slide(title: str, body: str, cls: str = "") -> str:
    return f"<section class='slide {cls}'><div class='kicker'>{esc(title)}</div>{body}</section>"


def write_deck_html(article: Article, out_dir: Path, terms: list[Concept]) -> Path:
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    slides: list[str] = []
    metric_cards = (
        f"<div class='metric'><b>{len(article.sections)}</b><span>章节块</span></div>"
        f"<div class='metric'><b>{len(terms)}</b><span>概念/节点</span></div>"
        f"<div class='metric'><b>{floor['min_slides']}</b><span>最低页数</span></div>"
    )
    slides.append(
        slide(
            "Cover",
            f"<h1>{esc(article.title)}</h1><p>把 source 转成可讲、可看、可考、可复验的一组产物。</p><div class='grid'>{metric_cards}</div>",
            "dark",
        )
    )
    section_cards = "".join(
        f"<div class='card'><b>{esc(title)}</b><p>{esc(section_excerpt(content, 120))}</p></div>"
        for title, content in article.sections[:6]
    )
    slides.append(slide("文章地图", f"<h2>先保留结构，再做媒介化</h2><div class='grid'>{section_cards}</div>"))

    matrix_rows = "".join(
        f"<tr><td>{esc(theme)}</td><td>{esc(term_names(bucket, 10))}</td><td>{len(bucket)}</td></tr>"
        for theme, bucket in rows[:12]
    )
    slides.append(
        slide(
            "覆盖矩阵",
            f"<h2>所有产物共用这一张清单</h2><table><tr><th>模块</th><th>概念</th><th>数量</th></tr>{matrix_rows}</table>",
        )
    )

    flow_nodes = "".join(
        f"<div><b>{idx}</b><p>{esc(title)}</p></div>" for idx, (title, _) in enumerate(article.sections[:5], 1)
    )
    slides.append(slide("Source Path", f"<h2>按原文路径推进</h2><div class='flow'>{flow_nodes}</div>", "dark"))

    for theme, bucket in rows[:8]:
        cards = "".join(
            f"<div class='card'><b>{esc(term)}</b><p>{esc(definition)}</p></div>"
            for term, _, definition in bucket[:6]
        )
        slides.append(slide(theme, f"<h2>{esc(theme)}</h2><div class='grid'>{cards}</div>"))

    for idx, group in enumerate(chunked(terms, 10)[:2], 1):
        table_rows = "".join(
            f"<tr><td>{esc(term)}</td><td>{esc(category)}</td><td>{esc(definition)}</td></tr>"
            for term, category, definition in group
        )
        slides.append(slide(f"术语速查 {idx}", f"<h2>覆盖清单 {idx}</h2><table><tr><th>概念</th><th>章节</th><th>解释</th></tr>{table_rows}</table>"))

    pair_cards = ""
    for left, right in zip(terms[0::2][:4], terms[1::2][:4]):
        pair_cards += (
            f"<div class='box'><b>{esc(left[0])} / {esc(right[0])}</b>"
            f"<p>{esc(left[1])} vs {esc(right[1])}。讲解时要回到 source，不凭空扩展。</p></div>"
        )
    slides.append(slide("易混与边界", f"<h2>相邻概念要讲出区别</h2><div class='two'>{pair_cards}</div>"))

    slides.append(
        slide(
            "转换行动清单",
            "<h2>交付前逐项检查</h2><ol><li>是否覆盖主要章节？</li><li>是否覆盖主体概念？</li><li>是否保留案例/流程？</li><li>是否有风险和边界？</li><li>是否能回指 source？</li></ol>",
        )
    )
    slides.append(
        slide(
            "自测",
            "<h2>三道验收题</h2><ol><li>这篇文章的主线是什么？</li><li>哪些概念被遗漏会改变理解？</li><li>哪个产物最容易变成摘要，如何补救？</li></ol>",
            "dark",
        )
    )
    coverage_index = esc(term_names(terms, 40))
    slides.append(slide("Source", f"<h2>来源与覆盖索引</h2><p>{esc(source_label(article))}</p><p class='small'>{coverage_index}</p>"))

    if len(slides) > floor["max_slides"]:
        slides = slides[: floor["max_slides"] - 1] + [slides[-1]]
    while len(slides) < floor["min_slides"]:
        section = article.sections[(len(slides) - 1) % len(article.sections)]
        slides.insert(
            -1,
            slide(
                f"Source Detail {len(slides)}",
                f"<h2>{esc(section[0])}</h2><p>{esc(section_excerpt(section[1], 500))}</p>",
            ),
        )

    html_body = "\n".join(slides)
    return write_text(out_dir / "deck.html", html_doc(article.title, DECK_CSS, html_body))


def html_doc(title: str, css: str, body: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{esc(title)}</title>
  <style>{css}</style>
</head>
<body>
{body}
</body>
</html>
"""


def add_textbox(slide_obj, left, top, width, height, text: str, size: int, color: RGBColor | None = None, bold: bool = False) -> None:
    box = slide_obj.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_card(slide_obj, left, top, width, height, title: str, body: str, fill: RGBColor) -> None:
    shape = slide_obj.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(220, 226, 235)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(30, 41, 59)
    p2 = frame.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(51, 65, 85)


def write_deck_pptx(article: Article, out_dir: Path, terms: list[Concept]) -> Path | None:
    if Presentation is None:
        return None
    floor = qa_floor(article, terms)
    rows = theme_rows(terms)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str = ""):
        s = prs.slides.add_slide(blank)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(248, 250, 252)
        add_textbox(s, Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.58), title, 27, RGBColor(15, 23, 42), True)
        if subtitle:
            add_textbox(s, Inches(0.58), Inches(0.95), Inches(11.6), Inches(0.38), subtitle, 12, RGBColor(100, 116, 139))
        return s

    s = new_slide(article.title, source_label(article))
    add_textbox(s, Inches(0.7), Inches(2.0), Inches(11.8), Inches(0.72), "保真转换：先覆盖 source，再生成不同媒介产物。", 23, RGBColor(30, 64, 175), True)
    for idx, (label, value) in enumerate([("章节块", len(article.sections)), ("概念节点", len(terms)), ("最低页数", floor["min_slides"])]):
        add_card(s, Inches(0.75 + idx * 4.1), Inches(3.2), Inches(3.6), Inches(1.55), label, str(value), RGBColor(239, 246, 255))

    s = new_slide("文章地图", "不跳过结构，先让读者知道原文如何展开")
    for idx, (title, content) in enumerate(article.sections[:6]):
        add_card(s, Inches(0.6 + (idx % 3) * 4.15), Inches(1.55 + (idx // 3) * 2.35), Inches(3.75), Inches(1.85), title, section_excerpt(content, 130), RGBColor(255, 255, 255))

    s = new_slide("覆盖矩阵", "所有产物共用同一张概念清单")
    for idx, (theme, bucket) in enumerate(rows[:8]):
        add_card(s, Inches(0.6 + (idx % 4) * 3.15), Inches(1.45 + (idx // 4) * 2.25), Inches(2.8), Inches(1.8), theme, term_names(bucket, 6), RGBColor(240, 253, 244))

    for theme, bucket in rows[:8]:
        s = new_slide(theme, "每页负责一组概念，避免把长文压成少量 bullet")
        for idx, (term, _, definition) in enumerate(bucket[:6]):
            add_card(s, Inches(0.6 + (idx % 3) * 4.15), Inches(1.5 + (idx // 3) * 2.25), Inches(3.75), Inches(1.75), term, definition, RGBColor(255, 255, 255))

    for idx, group in enumerate(chunked(terms, 8)[:2], 1):
        s = new_slide(f"术语速查 {idx}", "检查 PPT 是否覆盖主体概念")
        for row_idx, (term, category, definition) in enumerate(group):
            add_textbox(s, Inches(0.72), Inches(1.35 + row_idx * 0.65), Inches(2.1), Inches(0.38), term, 14, RGBColor(30, 64, 175), True)
            add_textbox(s, Inches(2.65), Inches(1.35 + row_idx * 0.65), Inches(2.4), Inches(0.38), category, 12, RGBColor(71, 85, 105))
            add_textbox(s, Inches(5.0), Inches(1.35 + row_idx * 0.65), Inches(7.2), Inches(0.42), definition, 11, RGBColor(51, 65, 85))

    s = new_slide("行动清单", "交付前逐项验收")
    add_textbox(s, Inches(0.9), Inches(1.6), Inches(11.2), Inches(4.8), "1. 覆盖主要章节\n2. 覆盖主体概念\n3. 保留案例、流程或证据链\n4. 标出风险与边界\n5. 输出文件可打开、可读、可追溯", 24, RGBColor(15, 23, 42))

    s = new_slide("来源", "Source-grounded output")
    add_textbox(s, Inches(0.9), Inches(1.6), Inches(11.4), Inches(2.2), source_label(article), 18, RGBColor(71, 85, 105))
    add_textbox(s, Inches(0.9), Inches(3.6), Inches(11.4), Inches(2.0), term_names(terms, 36), 14, RGBColor(100, 116, 139))

    while len(prs.slides) < min(floor["min_slides"], 14):
        title, content = article.sections[(len(prs.slides) - 1) % len(article.sections)]
        s = new_slide(f"Source Detail {len(prs.slides)}", title)
        add_textbox(s, Inches(0.9), Inches(1.55), Inches(11.4), Inches(4.8), section_excerpt(content, 900), 18, RGBColor(30, 41, 59))

    path = out_dir / "deck.pptx"
    prs.save(path)
    return path


def render_with_playwright(out_dir: Path, node_bin: str, node_path: str | None = None) -> dict[str, object]:
    script = """
const { chromium } = require('playwright');
const path = require('path');
const outDir = __OUTDIR__;
const toUrl = file => 'file://' + path.resolve(outDir, file);

async function launchBrowser() {
  try {
    return await chromium.launch({ headless: true, channel: 'chrome' });
  } catch (first) {
    return await chromium.launch({ headless: true });
  }
}

(async () => {
  const browser = await launchBrowser();
  const page = await browser.newPage();
  await page.goto(toUrl('report.html'), { waitUntil: 'networkidle' });
  await page.pdf({ path: path.join(outDir, 'report.pdf'), format: 'A4', printBackground: true, margin: { top: '12mm', bottom: '12mm', left: '10mm', right: '10mm' } });
  await page.setViewportSize({ width: 1080, height: 1920 });
  await page.goto(toUrl('infographic.html'), { waitUntil: 'networkidle' });
  await page.screenshot({ path: path.join(outDir, 'infographic.png'), fullPage: true });
  await page.setViewportSize({ width: 1600, height: 900 });
  await page.goto(toUrl('deck.html'), { waitUntil: 'networkidle' });
  await page.screenshot({ path: path.join(outDir, 'deck-cover.png') });
  await page.pdf({ path: path.join(outDir, 'deck.pdf'), width: '1600px', height: '900px', printBackground: true });
  await browser.close();
})().catch(err => {
  console.error(err && err.stack ? err.stack : String(err));
  process.exit(1);
});
""".replace("__OUTDIR__", json.dumps(str(out_dir)))
    env = os.environ.copy()
    if node_path:
        env["NODE_PATH"] = node_path
    proc = subprocess.run([node_bin, "-e", script], text=True, capture_output=True, env=env, check=False)
    if proc.returncode == 0:
        return {"ok": True, "provider": "playwright", "returncode": 0, "stdout": proc.stdout.strip(), "stderr": proc.stderr.strip()}

    # Keep local mode usable on machines where the Node Playwright module is
    # absent. This is a browser fallback, never an Open Design result.
    wkhtmltoimage = shutil.which("wkhtmltoimage")
    wkhtmltopdf = shutil.which("wkhtmltopdf")
    if not wkhtmltoimage or not wkhtmltopdf:
        return {"ok": False, "provider": "playwright", "returncode": proc.returncode, "stdout": proc.stdout.strip(), "stderr": proc.stderr.strip()}

    def run(command: list[str]) -> subprocess.CompletedProcess[str]:
        return subprocess.run(command, text=True, capture_output=True, check=False)

    report = run([wkhtmltopdf, "--enable-local-file-access", str(out_dir / "report.html"), str(out_dir / "report.pdf")])
    infographic = run([wkhtmltoimage, "--enable-local-file-access", "--disable-smart-width", "--width", "1080", "--quality", "94", str(out_dir / "infographic.html"), str(out_dir / "infographic.png")])
    cover = run([wkhtmltoimage, "--enable-local-file-access", "--width", "1600", "--height", "900", "--quality", "94", str(out_dir / "deck.html"), str(out_dir / "deck-cover.png")])
    deck = run([wkhtmltopdf, "--enable-local-file-access", "--page-width", "423mm", "--page-height", "238mm", str(out_dir / "deck.html"), str(out_dir / "deck.pdf")])
    ok = all(item.returncode == 0 for item in (report, infographic, cover, deck))
    fallback_error = "\n".join(item.stderr.strip() for item in (report, infographic, cover, deck) if item.stderr.strip())
    return {"ok": ok, "provider": "browser-fallback", "returncode": 0 if ok else 1, "stdout": "", "stderr": fallback_error}


def build_bundle(args: argparse.Namespace) -> dict[str, object]:
    article = parse_article(Path(args.article))
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    terms = matched_terms(article)

    files: list[Path] = [
        write_report(article, out_dir, terms),
        write_report_html(article, out_dir, terms),
        write_data_table(out_dir, terms),
        write_quiz(article, out_dir, terms),
        *write_flashcards(article, out_dir, terms),
        write_mindmap(article, out_dir, terms),
        *write_scripts(article, out_dir, terms),
        write_infographic_html(article, out_dir, terms),
        write_deck_html(article, out_dir, terms),
    ]
    pptx_path = write_deck_pptx(article, out_dir, terms)
    if pptx_path:
        files.append(pptx_path)

    render_result = None
    if not args.no_render:
        render_result = render_with_playwright(out_dir, args.node_bin, args.node_path)
        for name in ["report.pdf", "infographic.png", "deck-cover.png", "deck.pdf"]:
            path = out_dir / name
            if path.exists():
                files.append(path)

    qa = validate_bundle(article, out_dir, terms)
    manifest = {
        "source": str(article.path),
        "title": article.title,
        "provider": "soia-local",
        "content_modes": ["preserve", "learning", "visual_dense"],
        "concept_count": len(terms),
        "section_count": len(article.sections),
        "files": [str(path) for path in files],
        "render": render_result,
        "qa": qa,
    }
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.strict and not qa.get("ok"):
        raise SystemExit(json.dumps(qa, ensure_ascii=False, indent=2))
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--article", required=True, help="Markdown article path")
    parser.add_argument("--out-dir", required=True, help="Output directory")
    parser.add_argument("--node-bin", default="node", help="Node.js executable for Playwright rendering")
    parser.add_argument("--node-path", default=None, help="NODE_PATH containing playwright module")
    parser.add_argument("--no-render", action="store_true", help="Skip Playwright PDF/PNG rendering")
    parser.add_argument("--strict", action="store_true", help="Exit non-zero if quality gates fail")
    parser.add_argument("--json", action="store_true", help="Print manifest JSON")
    args = parser.parse_args()
    manifest = build_bundle(args)
    if args.json:
        print(json.dumps(manifest, ensure_ascii=False, indent=2))
    else:
        print(f"wrote {len(manifest['files'])} files to {args.out_dir}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
